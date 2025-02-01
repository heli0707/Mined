from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class PPTImageAppender:
    def __init__(self, ppt_file, image_folder):
        self.ppt_file = ppt_file
        self.image_folder = image_folder
        self.prs = None

    def load_ppt(self):
        """Load the PowerPoint presentation."""
        if not os.path.exists(self.ppt_file):
            print(f"Error: The file {self.ppt_file} does not exist.")
            return False
        
        try:
            self.prs = Presentation(self.ppt_file)
            print(f"Loaded presentation {self.ppt_file}.")
            return True
        except Exception as e:
            print(f"Error loading PowerPoint file: {e}")
            return False

    def add_images(self):
        """Append images from the folder to the PowerPoint."""
        if not self.prs:
            print("Error: No PowerPoint file loaded.")
            return
        
        if not os.path.isdir(self.image_folder):
            print(f"Error: The folder {self.image_folder} does not exist.")
            return
        
        image_files = [f for f in os.listdir(self.image_folder) if f.endswith('.png')]
        
        if not image_files:
            print(f"No PNG images found in {self.image_folder}.")
            return
        
        image_files.sort()
        
        # Iterate through all image files
        for image_file in image_files:
            image_path = os.path.join(self.image_folder, image_file)
            
            slide_layout = self.prs.slide_layouts[5]  # Blank slide
            slide = self.prs.slides.add_slide(slide_layout)
            
            left = Inches(1)
            top = Inches(1)
            width = Inches(6)
            height = Inches(4)
            
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"Image {image_file} added to slide.")

    def add_thank_you_slide(self):
        """Add a 'Thank You' slide at the end."""
        if not self.prs:
            print("Error: No PowerPoint file loaded.")
            return
        
        thank_you_layout = self.prs.slide_layouts[6]  # Blank slide
        thank_you_slide = self.prs.slides.add_slide(thank_you_layout)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(2)
        
        txBox = thank_you_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        p = tf.add_paragraph()
        p.text = "THANK YOU"
        p.alignment = PP_ALIGN.CENTER
        
        font = p.font
        font.size = Pt(72)  # Large font size
        font.bold = True
        font.color.rgb = RGBColor(44, 116, 181)  # Blue color
        
        print("Thank You slide added.")

    def save_ppt(self):
        """Save the modified PowerPoint file."""
        if not self.prs:
            print("Error: No PowerPoint file loaded.")
            return
        
        try:
            self.prs.save(self.ppt_file)
            print(f"PowerPoint saved as '{self.ppt_file}'.")
        except Exception as e:
            print(f"Error saving PowerPoint: {e}")

    def append_images_and_add_thank_you(self):
        """Perform all steps to append images and add 'Thank You' slide."""
        if not self.load_ppt():
            return
        
        self.add_images()
        self.add_thank_you_slide()
        self.save_ppt()

# Example usage
ppt_file_path = "final_effects.pptx"
image_folder_path = "images"

ppt_appender = PPTImageAppender(ppt_file_path, image_folder_path)
ppt_appender.append_images_and_add_thank_you()
