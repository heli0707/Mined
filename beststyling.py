# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_SHAPE

# class PPTStyleApplicator:
#     THEME_COLORS = {
#         'primary': RGBColor(0, 75, 135),     # Dark Blue
#         'secondary': RGBColor(0, 130, 200),   # Medium Blue
#         'accent': RGBColor(245, 245, 245),    # Light Gray
#         'text': RGBColor(51, 51, 51)         # Dark Gray
#     }
#     FONT_NAME = "Arial"
#     TITLE_FONT_SIZE = Pt(40)
#     SUBTITLE_FONT_SIZE = Pt(28)
#     HEADING_FONT_SIZE = Pt(24)
#     BODY_FONT_SIZE = Pt(15)

#     def _init_(self):
#         # Define text formatting
#         # self.FONT_NAME = "Arial"
#         # self.TITLE_FONT_SIZE = Pt(40)
#         # self.SUBTITLE_FONT_SIZE = Pt(28)
#         # self.HEADING_FONT_SIZE = Pt(24)
#         # self.BODY_FONT_SIZE = Pt(15)
        
#         # Theme colors
#         self.THEME_COLORS = {
#             'primary': RGBColor(0, 75, 135),     # Dark Blue
#             'secondary': RGBColor(0, 130, 200),   # Medium Blue
#             'accent': RGBColor(245, 245, 245),    # Light Gray
#             'text': RGBColor(51, 51, 51)         # Dark Gray
#         }

#     def apply_slide_template(self, slide):
#         """Applies the professional template to a slide."""
#         # Set background
#         background = slide.background
#         fill = background.fill
#         fill.solid()
#         fill.fore_color.rgb = self.THEME_COLORS['accent']
        
#         # Add bottom accent bar
#         left = 0
#         top = Inches(6.5)
#         width = Inches(13.333)
#         height = Inches(1)
        
#         shape = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE, 
#             left, 
#             top, 
#             width, 
#             height
#         )
#         fill = shape.fill
#         fill.solid()
#         fill.fore_color.rgb = self.THEME_COLORS['secondary']
#         shape.line.color.rgb = self.THEME_COLORS['secondary']

#     def style_title_slide(self, slide):
#         """Styles the title slide."""
#         if not slide.shapes.title:
#             return
            
#         title_shape = slide.shapes.title
#         title_frame = title_shape.text_frame
        
#         # Style the title
#         title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
#         paragraph = title_frame.paragraphs[0]
#         paragraph.font.size = self.TITLE_FONT_SIZE
#         paragraph.font.name = self.FONT_NAME
#         paragraph.font.color.rgb = self.THEME_COLORS['primary']
#         paragraph.alignment = PP_ALIGN.CENTER
        
#         # Add centered decorative line
#         left = Inches(4)
#         top = Inches(4.5)
#         width = Inches(2)
#         height = Inches(0.1)
#         shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#         fill = shape.fill
#         fill.solid()
#         fill.fore_color.rgb = self.THEME_COLORS['secondary']

#     def style_content_slide(self, slide):
#         """Styles a content slide."""
#         # Style title if exists
#         if slide.shapes.title:
#             title_shape = slide.shapes.title
#             title_frame = title_shape.text_frame
#             paragraph = title_frame.paragraphs[0]
#             paragraph.font.size = self.HEADING_FONT_SIZE
#             paragraph.font.name = self.FONT_NAME
#             paragraph.font.color.rgb = self.THEME_COLORS['primary']
#             paragraph.alignment = PP_ALIGN.CENTER

#         # Style content
#         for shape in slide.shapes:
#             if not hasattr(shape, "text_frame"):
#                 continue
                
#             text_frame = shape.text_frame
#             text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
#             for paragraph in text_frame.paragraphs:
#                 if paragraph.text.strip():
#                     paragraph.font.size = self.BODY_FONT_SIZE
#                     paragraph.font.name = self.FONT_NAME
#                     paragraph.font.color.rgb = self.THEME_COLORS['text']
#                     paragraph.alignment = PP_ALIGN.CENTER
#                     paragraph.space_before = Pt(6)
#                     paragraph.space_after = Pt(6)
#                     paragraph.line_spacing = 1.2

#     def apply_styling(self, input_path, output_path):
#         """Applies styling to an existing PowerPoint presentation."""
#         # Load the presentation
#         prs = Presentation(input_path)
        
#         # Set slide dimensions to widescreen
#         prs.slide_width = Inches(13.333)
#         prs.slide_height = Inches(7.5)
        
#         # Process each slide
#         for i, slide in enumerate(prs.slides):
#             # Apply base template
#             self.apply_slide_template(slide)
            
#             # Apply specific styling based on slide type
#             if i == 0:  # First slide (title slide)
#                 self.style_title_slide(slide)
#             else:  # Content slides
#                 self.style_content_slide(slide)
        
#         # Save the styled presentation
#         prs.save(output_path)
#         print(f"✅ Successfully applied styling to {output_path}")

# def main():
#     # Example usage
#     applicator = PPTStyleApplicator()
#     applicator.apply_styling(
#         input_path="final_effects.pptx",
#         output_path="final_effects.pptx"
#     )

# if __name__ == "__main__":
#     main()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import logging

class PPTStyleApplicator:
    def __init__(self):
        # Set up logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)

        # Define theme colors
        self.THEME_COLORS = {
            'primary': RGBColor(0, 75, 135),     # Dark Blue
            'secondary': RGBColor(0, 130, 200),   # Medium Blue
            'accent': RGBColor(245, 245, 245),    # Light Gray
            'text': RGBColor(51, 51, 51)         # Dark Gray
        }

        # Define text formatting
        self.FONT_NAME = "Arial"
        self.TITLE_FONT_SIZE = Pt(40)
        self.SUBTITLE_FONT_SIZE = Pt(28)
        self.HEADING_FONT_SIZE = Pt(24)
        self.BODY_FONT_SIZE = Pt(15)

    def apply_slide_template(self, slide):
        """Applies the professional template to a slide."""
        try:
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.THEME_COLORS['accent']
            
            # Add bottom accent bar
            left = 0
            top = Inches(6.5)
            width = Inches(13.333)
            height = Inches(1)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                left, 
                top, 
                width, 
                height
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.THEME_COLORS['secondary']
            shape.line.color.rgb = self.THEME_COLORS['secondary']
        except Exception as e:
            self.logger.error(f"Error applying slide template: {str(e)}")

    def style_title_slide(self, slide):
        """Styles the title slide."""
        try:
            if not slide.shapes.title:
                return
                
            title_shape = slide.shapes.title
            title_frame = title_shape.text_frame
            
            # Style the title
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = self.TITLE_FONT_SIZE
            paragraph.font.name = self.FONT_NAME
            paragraph.font.color.rgb = self.THEME_COLORS['primary']
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Add centered decorative line
            left = Inches(4)
            top = Inches(4.5)
            width = Inches(2)
            height = Inches(0.1)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.THEME_COLORS['secondary']
        except Exception as e:
            self.logger.error(f"Error styling title slide: {str(e)}")

    def style_content_slide(self, slide):
        """Styles a content slide."""
        try:
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_frame = title_shape.text_frame
                paragraph = title_frame.paragraphs[0]
                paragraph.font.size = self.HEADING_FONT_SIZE
                paragraph.font.name = self.FONT_NAME
                paragraph.font.color.rgb = self.THEME_COLORS['primary']
                paragraph.alignment = PP_ALIGN.CENTER

            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                    
                text_frame = shape.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraph.font.size = self.BODY_FONT_SIZE
                        paragraph.font.name = self.FONT_NAME
                        paragraph.font.color.rgb = self.THEME_COLORS['text']
                        paragraph.alignment = PP_ALIGN.CENTER
                        paragraph.space_before = Pt(6)
                        paragraph.space_after = Pt(6)
                        paragraph.line_spacing = 1.2
        except Exception as e:
            self.logger.error(f"Error styling content slide: {str(e)}")

    def apply_styling(self, input_path, output_path):
        """Applies styling to an existing PowerPoint presentation."""
        try:
            # Load the presentation
            prs = Presentation(input_path)
            
            # Set slide dimensions to widescreen
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Process each slide
            for i, slide in enumerate(prs.slides):
                # Apply base template
                self.apply_slide_template(slide)
                
                # Apply specific styling based on slide type
                if i == 0:  # First slide (title slide)
                    self.style_title_slide(slide)
                else:  # Content slides
                    self.style_content_slide(slide)
            
            # Save the styled presentation
            prs.save(output_path)
            self.logger.info(f"✅ Successfully applied styling to {output_path}")
            
        except Exception as e:
            self.logger.error(f"applied styling: {str(e)}")
            return 0