import os
import fitz
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import openai
import textwrap
import re
from dotenv import load_dotenv
import os
import groq

# Load environment variables & Retrieve the Groq API key from .env
load_dotenv()
api_key = os.getenv("GROQ_API_KEY")

if not api_key:
    print("Error: GROQ_API_KEY not found. Make sure it's set in the .env file.")
else:
    client = groq.Groq(api_key=api_key)


class PDFtoPPTGenerator:
    def __init__(self):
        logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
        self.logger = logging.getLogger(__name__)

        # Define text formatting
        self.FONT_NAME = "Arial"
        self.TITLE_FONT_SIZE = Pt(40)
        self.SUBTITLE_FONT_SIZE = Pt(28)
        self.HEADING_FONT_SIZE = Pt(24)
        self.BODY_FONT_SIZE = Pt(15)

        # Load environment variables
        load_dotenv()
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
            raise ValueError("GROQ_API_KEY not found in environment variables")
        self.client = groq.Groq(api_key=api_key)
        
        # Theme colors
        self.THEME_COLORS = {
            'primary': RGBColor(0, 75, 135),     # Dark Blue
            'secondary': RGBColor(0, 130, 200),   # Medium Blue
            'accent': RGBColor(245, 245, 245),    # Light Gray
            'text': RGBColor(51, 51, 51)         # Dark Gray
        }
        
        self.MAX_BULLETS_PER_SLIDE = 4
        self.logger.info("✅ PDF to PPT Generator Initialized!")
    def apply_slide_template(self, slide):
      """Applies a professional template to the slide."""
      background = slide.background
      fill = background.fill
      fill.solid()
      fill.fore_color.rgb = self.THEME_COLORS['accent']
      
      # Add a full-width accent bar at the bottom
      left = 0  # Start from the leftmost edge
      top = Inches(6.5)
      # Get presentation width from the slide's parent presentation
      width = Inches(13.333)  # Standard widescreen presentation width
      height = Inches(1)
      
      shape = slide.shapes.add_shape(
          MSO_SHAPE.RECTANGLE, 
          left, 
          top, 
          width, 
          height
      )
      fill = shape.fill
      fill.solid()
      fill.fore_color.rgb = self.THEME_COLORS['secondary']
      shape.line.color.rgb = self.THEME_COLORS['secondary']
    def create_title_slide(self, prs, title: str):
        """Creates a styled title slide with centered content."""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.apply_slide_template(slide)
        
        # Add main title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        
        # Center align the text vertically and horizontally
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Style the title
        paragraph = title_frame.paragraphs[0]
        paragraph.font.size = self.TITLE_FONT_SIZE
        paragraph.font.name = self.FONT_NAME
        paragraph.font.color.rgb = self.THEME_COLORS['primary']
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Remove subtitle placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Add centered decorative line
        left = Inches(4)
        top = Inches(4.5)
        width = Inches(2)
        height = Inches(0.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = self.THEME_COLORS['secondary']

    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """Extracts text from the entire PDF."""
        try:
            doc = fitz.open(pdf_path)
            text = "\n".join([page.get_text("text") for page in doc])
            return text.strip() if text else ""
        except Exception as e:
            self.logger.error(f"❌ Error extracting text from PDF: {str(e)}")
            return ""

    def extract_references(self, pdf_path: str) -> list:
        """Extracts the references section from the PDF and returns it as a list of lines."""
        doc = fitz.open(pdf_path)
        references = []
        capture = False

        for page in doc:
            text = page.get_text("text")
            lines = text.split("\n")

            for line in lines:
                if "references" in line.lower():
                    capture = True
                    continue

                if capture:
                    references.append(line.strip())

        return references if references else ["No references available."]

    def extract_title(self, text: str) -> str:
        """Extracts the research paper title (assumed to be the first non-empty line)."""
        lines = text.split("\n")
        for line in lines:
            if line.strip():
                return line.strip()
        return "Research Paper Title"

    def detect_paper_type(self, text: str) -> str:
        """Uses OpenAI Mixtral to classify the research paper type."""
        try:
            response = client.chat.completions.create(
                model="mixtral-8x7b-32768",
                messages=[
                    {"role": "system", "content": "You are an AI expert in research paper classification."},
                    {"role": "user", "content": f"Analyze this research paper and classify it into one of the following categories:\n\n"
                                              "- Computer Science Implementation Paper\n"
                                              "- Computer Science Review Paper\n"
                                              "- Physics Paper\n"
                                              "- Chemistry Paper\n"
                                              "- Biology Paper\n"
                                              "- Mathematics Paper\n"
                                              "- Electrical Engineering Paper\n"
                                              "- Mechanical Engineering Paper\n\n"
                                              "Provide only the category name."},
                ],
                max_tokens=50
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            self.logger.error(f"❌ Error detecting paper type: {str(e)}")
            return "Unknown Paper Type"

    def split_text_into_chunks(self, text: str) -> list:
        """Splits text into chunks that fit on a slide."""
        if not text:
            return ["Content unavailable."]
            
        # Split text into sentences first
        sentences = text.replace('\n', ' ').split('. ')
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            # Add period back if it was removed during split
            if sentence and not sentence.endswith('.'):
                sentence += '.'
                
            # If adding this sentence would exceed the limit, start a new chunk
            if len(current_chunk) + len(sentence) > self.MAX_CHARS_PER_SLIDE:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + ' '
            else:
                current_chunk += sentence + ' '
                
        # Add the last chunk if it's not empty
        if current_chunk:
            chunks.append(current_chunk.strip())
            
        return chunks if chunks else ["Content unavailable."]

    def convert_to_bullets(self, text: str) -> list:
        """Converts text into bullet points, ensuring sentences stay together."""
        if not text:
            return ["Content unavailable."]
        
        # Clean and split the text into sentences
        text = text.replace('\n', ' ').strip()
        sentences = re.split('(?<=[.!?])\s+', text)
        
        # Group sentences into meaningful bullet points
        bullets = []
        current_bullet = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            # Start a new bullet if the current one would be too long
            if len(current_bullet) + len(sentence) > 150 and current_bullet:  # 150 chars max per bullet
                bullets.append(current_bullet.strip())
                current_bullet = sentence
            else:
                if current_bullet:
                    current_bullet += " " + sentence
                else:
                    current_bullet = sentence
        
        # Add the last bullet if not empty
        if current_bullet:
            bullets.append(current_bullet.strip())
        
        return bullets if bullets else ["Content unavailable."]

    def split_bullets_into_slides(self, bullets: list) -> list:
        """Splits bullets into slide-sized chunks."""
        return [bullets[i:i + self.MAX_BULLETS_PER_SLIDE] 
                for i in range(0, len(bullets), self.MAX_BULLETS_PER_SLIDE)]

    def add_content_slide(self, prs, title: str, bullets: list, slide_number: int = None):
        """Adds a styled content slide with centered content."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        self.apply_slide_template(slide)
        
        # Style the title
        title_shape = slide.shapes.title
        title_shape.text = title if slide_number is None else f"{title} (Continued {slide_number})"
        title_frame = title_shape.text_frame
        paragraph = title_frame.paragraphs[0]
        paragraph.font.size = self.HEADING_FONT_SIZE
        paragraph.font.name = self.FONT_NAME
        paragraph.font.color.rgb = self.THEME_COLORS['primary']
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Add content with center alignment
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.clear()
        
        for bullet in bullets:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = self.BODY_FONT_SIZE
            paragraph.font.name = self.FONT_NAME
            paragraph.font.color.rgb = self.THEME_COLORS['text']
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Adjust spacing
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.2

    def create_ppt(self, content: dict, output_path: str):
        """Creates a styled PowerPoint presentation."""
        prs = Presentation()
        
        # Set slide dimensions to widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Create slides with styling
        self.create_title_slide(prs, content.get("title", "Research Paper Title"))

        # Add content slides
        for section, text in content.items():
            if section not in ["title", "references"]:
                if isinstance(text, str):
                    bullets = self.convert_to_bullets(text)
                    bullet_slides = self.split_bullets_into_slides(bullets)
                    
                    for i, slide_bullets in enumerate(bullet_slides, 1):
                        slide_number = i if len(bullet_slides) > 1 else None
                        self.add_content_slide(prs, section.replace("_", " ").title(), 
                                            slide_bullets, slide_number)

        # Add styled references slides
        if "references" in content:
            references = content["references"]
            ref_chunks = [references[i:i + 5] for i in range(0, len(references), 5)]
            
            for i, chunk in enumerate(ref_chunks, 1):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                self.apply_slide_template(slide)
                
                # Style references title
                title_shape = slide.shapes.title
                title_shape.text = f"References (Page {i})"
                title_frame = title_shape.text_frame
                paragraph = title_frame.paragraphs[0]
                paragraph.font.size = self.HEADING_FONT_SIZE
                paragraph.font.name = self.FONT_NAME
                paragraph.font.color.rgb = self.THEME_COLORS['primary']
                
                # Add references with styling
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                
                for ref in chunk:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = ref
                    paragraph.level = 0
                    paragraph.font.size = self.BODY_FONT_SIZE
                    paragraph.font.name = self.FONT_NAME
                    paragraph.font.color.rgb = self.THEME_COLORS['text']

        # Save PowerPoint
        prs.save(output_path)
        print(f"✅ Successfully saved PowerPoint to {output_path}")



    def summarize_text(self, text: str, max_length: int = 512) -> str:
        """Summarizes text using OpenAI API."""
        try:
            response = client.chat.completions.create(
                model="mixtral-8x7b-32768",
                messages=[
                    {"role": "system", "content": "You are an AI assistant that summarizes research papers. Format your summary as clear, concise sentences that can be easily converted to presentation bullet points."},
                    {"role": "user", "content": f"Summarize the following research paper section:\n\n{text[:2000]}"}
                ],
                max_tokens=max_length
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            self.logger.error(f"❌ Error during summarization: {str(e)}")
            return "Summary not available."
    def process(self, pdf_path: str, output_path: str):
        """Processes the PDF and generates the PPT with dynamic file paths."""
        try:
            text = self.extract_text_from_pdf(pdf_path)
            if not text:
                self.logger.error("❌ No text extracted from PDF.")
                return

            title = self.extract_title(text)
            references = self.extract_references(pdf_path)
            paper_type = self.detect_paper_type(text)

            # Create sections based on paper type
            sections = self._create_sections(text, paper_type, references)
            content = {"title": f"{title} ({paper_type})", **sections}
            
            self.create_ppt(content, output_path)
            self.logger.info(f"✅ Successfully created presentation: {output_path}")
            
        except Exception as e:
            self.logger.error(f"❌ Error during PPT generation: {str(e)}")
            raise

    def _create_sections(self, text, paper_type, references):
        """Helper method to create sections based on paper type."""
        if "Review Paper" in paper_type:
            return {
                "Abstract": self.summarize_text(text[:2000]),
                "Introduction": self.summarize_text(text[2000:4000]),
                "Related Work": self.summarize_text(text[4000:6000]),
                "Discussion": self.summarize_text(text[6000:8000]),
                "Conclusion": self.summarize_text(text[8000:10000]),
                "References": references,
            }
        elif "Implementation Paper" in paper_type:
            return {
                "Abstract": self.summarize_text(text[:2000]),
                "Introduction": self.summarize_text(text[2000:4000]),
                "Methodology": self.summarize_text(text[4000:6000]),
                "Experiments & Results": self.summarize_text(text[6000:8000]),
                "Conclusion & Future Work": self.summarize_text(text[8000:10000]),
                "References": references,
            }
        else:
            return {
                "Abstract": self.summarize_text(text[:2000]),
                "Introduction": self.summarize_text(text[2000:4000]),
                "Main Content": self.summarize_text(text[4000:8000]),
                "Conclusion": self.summarize_text(text[8000:10000]),
                "References": references,
            }